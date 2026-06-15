from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
BG = RGBColor(0xF5, 0xF5, 0xF5)
HEADER_BG = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)

# Status lane colors (vertical columns)
STATUS_COLORS = {
    'Draft': RGBColor(0x95, 0xA5, 0xA6),
    'Processing Shapefile': RGBColor(0x5D, 0xAD, 0xEC),
    'With Operator': RGBColor(0x2E, 0xCC, 0x71),
    'With Reviewer': RGBColor(0xF3, 0x9C, 0x12),
    'Review Completed': RGBColor(0x8E, 0x44, 0xAD),
}

# Role colors
ROLE_COLORS = {
    'User': RGBColor(0xBD, 0xC3, 0xC7),
    'Operator': RGBColor(0x2E, 0xCC, 0x71),
    'Reviewer': RGBColor(0xF3, 0x9C, 0x12),
    'Silrec Admin': RGBColor(0xE7, 0x4C, 0x3C),
}

# Statuses in order with their data
STATUSES = [
    ('Draft', ['Operator can upload & process shapefile']),
    ('Processing Shapefile', ['System processes shapefile', 'Operator / Silrec Admin can send to With Operator']),
    ('With Operator', ['Operator completes data entry (cohorts, treatments)', 'Operator can send to With Reviewer', 'Operator can return to Draft']),
    ('With Reviewer', ['Reviewer reviews and approves/rejects', 'Reviewer can send back to With Operator', 'Reviewer can move to Review Completed']),
    ('Review Completed', ['Final status - proposal closed', 'Only Silrec Admin can reopen to With Reviewer']),
]

# Role definitions
ROLES = [
    ('User', 'View-only access. Can see proposals but cannot change status.'),
    ('Operator', 'Can upload shapefiles, enter data, send to reviewer. Core data entry role.'),
    ('Reviewer', 'Can review, approve, or return proposals to operator.'),
    ('Silrec Admin', 'Full access. Can act as any role and override status transitions.'),
]

# Transition arrows data: (from_idx, to_idx, label, roles)
TRANSITIONS = [
    (0, 1, 'Upload & process\nshapefile', 'Operator'),
    (1, 2, 'Processing complete\n→ With Operator', 'Operator, Silrec Admin'),
    (2, 3, 'Send to\nReviewer', 'Operator, Silrec Admin'),
    (2, 0, 'Return to\nDraft', 'Operator, Silrec Admin'),
    (3, 4, 'Approve\n→ Review Completed', 'Reviewer, Silrec Admin'),
    (3, 2, 'Send back\nto Operator', 'Reviewer, Silrec Admin'),
    (4, 3, 'Reopen', 'Silrec Admin'),
]

def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_box(slide, left, top, width, height, fill_color, text, font_size=13, bold=False, font_color=WHITE, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_arrow(slide, start_left, start_top, end_left, end_top, color=GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_textbox(slide, left, top, width, height, text, font_size=12, font_color=BLACK, bold=False, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

# ================ SLIDE 1: TITLE ================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, HEADER_BG)
add_textbox(slide, Inches(1), Inches(2.5), Inches(14), Inches(1.5),
    'SILREC - Processing Status & Role Workflow', font_size=41, font_color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(4.5), Inches(14), Inches(1),
    'Proposal Lifecycle: Status transitions and the roles that perform them', font_size=21, font_color=RGBColor(0xBD, 0xC3, 0xC7))
add_textbox(slide, Inches(1), Inches(6), Inches(14), Inches(0.5),
    'silrec.dbca.wa.gov.au', font_size=16, font_color=RGBColor(0x7F, 0x8C, 0x8D))

# ================ SLIDE 2: MAIN SWIMLANE DIAGRAM ================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Title
add_textbox(slide, Inches(0.5), Inches(0.2), Inches(15), Inches(0.6),
    'SILREC Processing Status Workflow', font_size=25, font_color=BLACK, bold=True, alignment=PP_ALIGN.LEFT)

# Legend
legend_y = Inches(0.3)
roles_legend = ['User (view)', 'Operator', 'Reviewer', 'Silrec Admin']
role_cols = [ROLE_COLORS['User'], ROLE_COLORS['Operator'], ROLE_COLORS['Reviewer'], ROLE_COLORS['Silrec Admin']]
for i, (rl, rc) in enumerate(zip(roles_legend, role_cols)):
    x = Inches(9.5) + Inches(i * 1.6)
    add_rounded_box(slide, x, legend_y, Inches(0.25), Inches(0.25), rc, '', font_size=7)
    add_textbox(slide, x + Inches(0.3), legend_y - Inches(0.05), Inches(1.3), Inches(0.35),
        rl, font_size=10, font_color=BLACK, alignment=PP_ALIGN.LEFT)

# Vertical lanes
num_statuses = len(STATUSES)
lane_width = Inches(2.6)
lane_gap = Inches(0.35)
start_x = Inches(0.5)
lane_top = Inches(1.0)
lane_height = Inches(5.5)

# Column header area
header_height = Inches(0.8)
role_area_top = lane_top + header_height + Inches(0.15)
role_area_height = Inches(3.2)
action_area_top = role_area_top + role_area_height + Inches(0.2)
action_area_height = Inches(1.5)

# Store center X positions for each column for arrow drawing
col_centers = []

for idx, (status_name, actions) in enumerate(STATUSES):
    x = start_x + idx * (lane_width + lane_gap)
    center_x = x + lane_width / 2
    col_centers.append((center_x, x))

    # Column background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, lane_top, lane_width, lane_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    shape.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
    shape.line.width = Pt(0.5)

    # Status header
    color = STATUS_COLORS[status_name]
    add_rounded_box(slide, x, lane_top, lane_width, header_height, color,
        status_name, font_size=15, bold=True)

    # Roles section label
    add_textbox(slide, x + Inches(0.1), role_area_top, lane_width - Inches(0.2), Inches(0.25),
        'Roles', font_size=9, font_color=GRAY, bold=True, alignment=PP_ALIGN.LEFT)

    # Actions section label
    add_textbox(slide, x + Inches(0.1), action_area_top, lane_width - Inches(0.2), Inches(0.25),
        'Actions', font_size=9, font_color=GRAY, bold=True, alignment=PP_ALIGN.LEFT)

    # Role boxes inside the lane
    act_top = action_area_top + Inches(0.3)
    for act_text in actions:
        add_rounded_box(slide, x + Inches(0.15), act_top, lane_width - Inches(0.3), Inches(0.45),
            RGBColor(0xE8, 0xF0, 0xFE), act_text, font_size=9, font_color=BLACK, alignment=PP_ALIGN.LEFT)
        act_top += Inches(0.55)

# Draw transitions as arrows - find center X for each status column
col_centers_x = [start_x + i * (lane_width + lane_gap) + lane_width / 2 for i in range(num_statuses)]

for from_s, to_s, label, roles in TRANSITIONS:
    x1 = col_centers_x[from_s]
    x2 = col_centers_x[to_s]
    y_pos = lane_top + lane_height / 2 - Inches(0.3)  # slightly above center to leave room

    # Arrow line
    arrow_top = y_pos
    connector = slide.shapes.add_connector(1, Emu(int(x1 * 914400)), Emu(int(arrow_top * 914400)),
                                            Emu(int(x2 * 914400)), Emu(int(arrow_top * 914400)))
    connector.line.color.rgb = RGBColor(0x5D, 0xAD, 0xEC)
    connector.line.width = Pt(2.5)

    # Arrow head triangle (small triangle at end)
    if x2 > x1:
        tri_left = x2 - Inches(0.15)
    else:
        tri_left = x2
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Emu(int(tri_left * 914400)),
                                  Emu(int((arrow_top - Inches(0.1)) * 914400)),
                                  Inches(0.2), Inches(0.2))
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(0x5D, 0xAD, 0xEC)
    tri.line.fill.background()
    if x2 < x1:
        tri.rotation = 180.0

    # Label above arrow
    label_x = min(x1, x2) + abs(x2 - x1) / 2
    lbl = add_textbox(slide, label_x - Inches(0.9), arrow_top - Inches(0.45), Inches(1.8), Inches(0.4),
        label, font_size=9, font_color=RGBColor(0x2C, 0x3E, 0x50), bold=True, alignment=PP_ALIGN.CENTER)

    # Role badge below arrow
    add_rounded_box(slide, label_x - Inches(0.7), arrow_top + Inches(0.05), Inches(1.4), Inches(0.25),
        RGBColor(0xE8, 0xF0, 0xFE), roles, font_size=8, font_color=RGBColor(0x5D, 0xAD, 0xEC))

# ================ SLIDE 3: ROLE DEFINITIONS ================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    'SILREC Role Definitions', font_size=25, font_color=BLACK, bold=True, alignment=PP_ALIGN.LEFT)

role_y = Inches(1.5)
for role_name, role_desc in ROLES:
    color = ROLE_COLORS[role_name] if role_name in ROLE_COLORS else GRAY
    add_rounded_box(slide, Inches(0.5), role_y, Inches(1.6), Inches(0.5),
        color, role_name, font_size=16, bold=True)
    add_textbox(slide, Inches(2.3), role_y, Inches(13), Inches(0.5),
        role_desc, font_size=15, font_color=BLACK, alignment=PP_ALIGN.LEFT)
    role_y += Inches(1.0)

# Group membership note
add_rounded_box(slide, Inches(0.5), Inches(5.8), Inches(15), Inches(0.8),
    RGBColor(0xFF, 0xF3, 0xCD),
    'Note: A user can belong to multiple groups. E.g. an Operator who also reviews can be assigned both Operator and Reviewer.',
    font_size=13, font_color=BLACK)

# ================ SLIDE 4: TRANSITION MATRIX ================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_textbox(slide, Inches(0.5), Inches(0.3), Inches(15), Inches(0.6),
    'Status Transition Matrix', font_size=25, font_color=BLACK, bold=True, alignment=PP_ALIGN.LEFT)

table_left = Inches(1.0)
table_top = Inches(1.5)
cols = 5
rows = 6
table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, Inches(14), Inches(3.5))
table = table_shape.table

# Set column widths
col_widths = [Inches(2.5), Inches(2.5), Inches(3.0), Inches(3.0), Inches(3.0)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

headers = ['From / To', 'With Operator', 'With Reviewer', 'Review Completed', 'Draft']
matrix = [
    ['Draft', 'Operator\n(upload & process)', '', '', ''],
    ['Processing\nShapefile', 'Operator, Silrec Admin\n(processing complete)', '', '', ''],
    ['With Operator', '', 'Operator, Silrec Admin\n(send to reviewer)', '', 'Operator, Silrec Admin\n(return to draft)'],
    ['With Reviewer', 'Reviewer, Silrec Admin\n(send back)', '', 'Reviewer, Silrec Admin\n(approve)', ''],
    ['Review\nCompleted', '', 'Silrec Admin\n(reopen)', '', ''],
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = HEADER_BG

for row_idx, row_data in enumerate(matrix):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            if col_idx == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = HEADER_BG
            else:
                paragraph.font.color.rgb = BLACK
            paragraph.alignment = PP_ALIGN.CENTER
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Save
output_path = '/home/ubuntu/projects/silrec/workflow.pptx'
prs.save(output_path)
print(f'Saved to {output_path}')
